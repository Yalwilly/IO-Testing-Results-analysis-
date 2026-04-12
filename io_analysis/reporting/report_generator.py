"""Report generator for IO Testing Results Analysis.

Generates PowerPoint-compatible reports with plots, summary tables,
and analysis comments.
"""

import logging
from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd

from io_analysis.config import Config
from io_analysis.data.models import AnalysisResult

logger = logging.getLogger(__name__)

# Color palette
COLOR_PASS = RGBColor(0x2E, 0xCC, 0x71)
COLOR_FAIL = RGBColor(0xE7, 0x4C, 0x3C)
COLOR_WARN = RGBColor(0xE6, 0x7E, 0x22)
COLOR_TITLE = RGBColor(0x2C, 0x3E, 0x50)
COLOR_SUBTITLE = RGBColor(0x7F, 0x8C, 0x8D)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
COLOR_DARK = RGBColor(0x34, 0x49, 0x5E)


def _add_title_slide(prs: Presentation, config: Config):
    """Add a title slide to the presentation."""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = config.report.title
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITLE
    title.text_frame.paragraphs[0].font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = (
        f"{config.report.subtitle}\n"
        f"{config.report.author}\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = COLOR_SUBTITLE


def _add_summary_slide(prs: Presentation, result: AnalysisResult, config: Config):
    """Add an executive summary slide."""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                     Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE

    summary = result.overall_summary
    pass_rate = summary.get("overall_pass_rate", 0)

    # Pass rate indicator (large)
    status_color = COLOR_PASS if pass_rate == 100 else (
        COLOR_WARN if pass_rate >= 95 else COLOR_FAIL
    )

    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3),
                                      Inches(4), Inches(2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    p = tf2.add_paragraph()
    p.text = f"Overall Pass Rate"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_SUBTITLE

    p = tf2.add_paragraph()
    p.text = f"{pass_rate:.1f}%"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = status_color

    # Summary table content
    summary_items = [
        ("Total Measurements", str(summary.get("total_measurements", 0))),
        ("Total Parameters", str(summary.get("total_parameters", 0))),
        ("Flows Tested", str(summary.get("total_flows", 0))),
        ("Total Pass", str(summary.get("total_pass", 0))),
        ("Total Fail", str(summary.get("total_fail", 0))),
        ("Parameters All Pass",
         str(summary.get("parameters_all_pass", 0))),
        ("Parameters With Failures",
         str(summary.get("parameters_with_fails", 0))),
    ]

    # Summary table
    rows = len(summary_items) + 1
    cols = 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(5), Inches(1.3),
                                         Inches(7), Inches(3.5))
    table = table_shape.table

    # Header
    for j, header in enumerate(["Metric", "Value"]):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = COLOR_WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK

    # Data rows
    for i, (metric, value) in enumerate(summary_items, 1):
        table.cell(i, 0).text = metric
        table.cell(i, 1).text = value
        for j in range(cols):
            cell = table.cell(i, j)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY

    # Key findings
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.0),
                                      Inches(12), Inches(3))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True

    p = tf3.add_paragraph()
    p.text = "Key Findings:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.space_after = Pt(6)

    # Add top comments (max 5)
    for comment in result.comments[:5]:
        p = tf3.add_paragraph()
        p.text = f"• {comment}"
        p.font.size = Pt(10)
        p.space_after = Pt(3)


def _add_plot_slide(prs: Presentation, plot_path: Path, title: str,
                    comment: str = "", config: Optional[Config] = None):
    """Add a slide with a plot image and optional comment."""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.1),
                                     Inches(12.5), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE

    # Plot image - sized to fit slide
    if plot_path.exists():
        img_left = Inches(0.5)
        img_top = Inches(0.8)
        img_width = Inches(12)
        img_height = Inches(5.5)
        slide.shapes.add_picture(str(plot_path), img_left, img_top,
                                 img_width, img_height)

    # Comment at bottom
    if comment:
        txBox2 = slide.shapes.add_textbox(Inches(0.3), Inches(6.5),
                                          Inches(12.5), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p = tf2.add_paragraph()
        p.text = comment
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_SUBTITLE


def _add_results_table_slide(prs: Presentation, result: AnalysisResult,
                             flow_name: str, config: Config):
    """Add a detailed results table slide for a flow."""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.1),
                                     Inches(12.5), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Detailed Results – {flow_name}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE

    # Collect data
    headers = ["Parameter", "Unit", "Spec Min", "Spec Max", "Mean", "Std",
               "Min", "Max", "Pass", "Fail", "Pass%", "Cpk", "Status"]

    params_in_flow = [
        p for p in result.all_parameters
        if (flow_name, p) in result.parameter_stats
    ]

    if not params_in_flow:
        return

    rows_count = min(len(params_in_flow), config.report.max_params_per_slide * 2)
    rows = rows_count + 1  # +1 for header

    table_shape = slide.shapes.add_table(
        rows, len(headers),
        Inches(0.2), Inches(0.8),
        Inches(12.9), Inches(6.2)
    )
    table = table_shape.table

    # Set column widths
    col_widths = [1.2, 0.6, 0.8, 0.8, 0.9, 0.9, 0.9, 0.9, 0.6, 0.6,
                  0.7, 0.7, 1.0]
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK

    # Data rows
    for i, param in enumerate(params_in_flow[:rows_count]):
        key = (flow_name, param)
        stats = result.parameter_stats[key]

        row_data = [
            stats.parameter,
            stats.unit,
            f"{stats.spec_min:.3f}" if stats.spec_min is not None else "–",
            f"{stats.spec_max:.3f}" if stats.spec_max is not None else "–",
            f"{stats.mean:.4f}",
            f"{stats.std:.4f}",
            f"{stats.minimum:.4f}",
            f"{stats.maximum:.4f}",
            str(stats.pass_count),
            str(stats.fail_count),
            f"{stats.pass_rate:.1f}",
            f"{stats.cpk:.2f}" if stats.cpk is not None else "–",
            stats.status,
        ]

        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.alignment = PP_ALIGN.CENTER

            # Color-code status and fail columns
            if j == len(headers) - 1:  # Status column
                if "PASS" in val and "FAIL" not in val:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xD5, 0xF5, 0xE3)
                elif "FAIL" in val:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFA, 0xDB, 0xD8)
                elif "MARGINAL" in val:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFD, 0xEA, 0xCD)

        # Alternating row colors
        if i % 2 == 0:
            for j in range(len(headers) - 1):
                cell = table.cell(i + 1, j)
                if not cell.fill.type:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY


def _add_comments_slide(prs: Presentation, result: AnalysisResult,
                        config: Config):
    """Add an analysis comments slide."""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.1),
                                     Inches(12.5), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Analysis Comments & Recommendations"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE

    # Comments
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.9),
                                      Inches(12), Inches(6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, comment in enumerate(result.comments):
        p = tf2.add_paragraph()

        # Different styling for different comment types
        if comment.startswith("OVERALL:"):
            p.text = comment
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COLOR_TITLE
            p.space_after = Pt(12)
        elif comment.startswith("[Cross-Flow]"):
            p.text = f"⟷ {comment}"
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_DARK
            p.space_after = Pt(6)
        else:
            # Determine icon based on content
            if "ALL PASS" in comment or "Capable" in comment:
                icon = "✓"
                color = COLOR_PASS
            elif "FAIL" in comment or "Not Capable" in comment:
                icon = "✗"
                color = COLOR_FAIL
            else:
                icon = "•"
                color = COLOR_DARK

            p.text = f"{icon} {comment}"
            p.font.size = Pt(9)
            p.font.color.rgb = color
            p.space_after = Pt(4)


def generate_csv_report(result: AnalysisResult, config: Config) -> Path:
    """Generate a CSV summary report.

    Returns path to the saved CSV file.
    """
    rows = []
    for (flow_name, param), stats in sorted(result.parameter_stats.items()):
        rows.append({
            "Flow": flow_name,
            "Parameter": param,
            "Unit": stats.unit,
            "Spec_Min": stats.spec_min,
            "Spec_Max": stats.spec_max,
            "Mean": stats.mean,
            "Std": stats.std,
            "Min": stats.minimum,
            "Max": stats.maximum,
            "Median": stats.median,
            "Sample_Count": stats.count,
            "Pass_Count": stats.pass_count,
            "Fail_Count": stats.fail_count,
            "Pass_Rate_%": stats.pass_rate,
            "Cpk": stats.cpk,
            "Status": stats.status,
            "Comment": stats.generate_comment(config.cpk_threshold),
        })

    df = pd.DataFrame(rows)
    csv_path = config.output_path / "analysis_results.csv"
    df.to_csv(csv_path, index=False)
    logger.info(f"CSV report saved: {csv_path}")
    return csv_path


def generate_pptx_report(result: AnalysisResult, plot_paths: dict,
                         config: Config) -> Path:
    """Generate a PowerPoint report with plots and analysis.

    Args:
        result: Complete AnalysisResult.
        plot_paths: Dict of category -> list of plot file paths.
        config: Configuration.

    Returns:
        Path to the saved PPTX file.
    """
    prs = Presentation()

    # Set slide dimensions (widescreen 16:9)
    prs.slide_width = Inches(config.report.slide_width_inches)
    prs.slide_height = Inches(config.report.slide_height_inches)

    # 1. Title slide
    _add_title_slide(prs, config)

    # 2. Executive summary
    _add_summary_slide(prs, result, config)

    # 3. Pass/Fail summary plots
    for plot_path in plot_paths.get("pass_fail_summary", []):
        flow = plot_path.stem.replace("pass_fail_summary_", "")
        _add_plot_slide(
            prs, plot_path,
            f"Pass/Fail Summary – {flow}",
            f"Stacked bar chart showing pass/fail rates for each IO parameter in {flow}.",
            config,
        )

    # 4. Detailed results tables
    for flow_name in result.all_flows:
        _add_results_table_slide(prs, result, flow_name, config)

    # 5. Parameter vs Spec plots
    for plot_path in plot_paths.get("parameter_vs_spec", []):
        param = plot_path.stem.replace("param_vs_spec_", "")
        # Find relevant comment
        comment = ""
        for c in result.comments:
            if param in c and "Cross-Flow" not in c:
                comment = c
                break
        _add_plot_slide(prs, plot_path,
                        f"{param} – Measured vs Spec Limits",
                        comment, config)

    # 6. Distribution histograms
    for plot_path in plot_paths.get("distributions", []):
        param = plot_path.stem.replace("histogram_", "")
        _add_plot_slide(prs, plot_path,
                        f"{param} – Distribution",
                        f"Distribution of measured values with spec limit overlays.",
                        config)

    # 7. Cross-flow comparison
    for plot_path in plot_paths.get("cross_flow", []):
        cross_flow_comments = [
            c for c in result.comments if "[Cross-Flow]" in c
        ]
        _add_plot_slide(
            prs, plot_path,
            "Cross-Flow Comparison",
            " | ".join(cross_flow_comments[:3]) if cross_flow_comments else "",
            config,
        )

    # 8. Cpk summary
    for plot_path in plot_paths.get("cpk_summary", []):
        flow = plot_path.stem.replace("cpk_summary_", "")
        _add_plot_slide(prs, plot_path,
                        f"Process Capability (Cpk) – {flow}",
                        f"Cpk values for all parameters. Target: {config.cpk_threshold}.",
                        config)

    # 9. Per-DUT scatter plots
    for plot_path in plot_paths.get("scatter_dut", []):
        param = plot_path.stem.replace("scatter_dut_", "")
        _add_plot_slide(prs, plot_path,
                        f"{param} – Per-DUT Results",
                        f"Individual DUT measurements with pass/fail indicators.",
                        config)

    # 10. Analysis comments
    _add_comments_slide(prs, result, config)

    # Save
    pptx_path = config.output_path / "IO_Validation_Report.pptx"
    prs.save(str(pptx_path))
    logger.info(f"PowerPoint report saved: {pptx_path}")

    return pptx_path


def generate_report(result: AnalysisResult, plot_paths: dict,
                    config: Config) -> dict:
    """Generate all report outputs.

    Returns dict of report type -> file path.
    """
    reports = {}

    # CSV report
    reports["csv"] = generate_csv_report(result, config)

    # PowerPoint report
    reports["pptx"] = generate_pptx_report(result, plot_paths, config)

    logger.info(f"Reports generated in {config.output_path}")
    return reports
